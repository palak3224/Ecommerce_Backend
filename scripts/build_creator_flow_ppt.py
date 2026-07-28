"""Generate the AOIN Creator x Seller (Wishlink-style) flow presentation as a .pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Brand palette ----
ORANGE = RGBColor(0xF2, 0x63, 0x1F)
DARK = RGBColor(0x1F, 0x2A, 0x37)
GREY = RGBColor(0x5B, 0x66, 0x70)
LIGHT = RGBColor(0xFA, 0xFA, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xFB, 0xE9, 0xE0)  # light orange tint
GREEN = RGBColor(0x2E, 0x7D, 0x32)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def rrect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    return shp


def txt(slide, x, y, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Calibri", italic=False, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb


def bullets(slide, x, y, w, h, items, size=16, color=DARK, gap=6, bullet="•  "):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        # support (text, bold) tuples
        if isinstance(it, tuple):
            label, rest = it
            r = p.add_run(); r.text = bullet + label
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = bullet + it
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


def header(slide, kicker, title):
    rect(slide, 0, 0, SW, Inches(1.15), WHITE)
    rect(slide, 0, 0, Inches(0.22), Inches(1.15), ORANGE)
    txt(slide, Inches(0.55), Inches(0.16), Inches(11), Inches(0.35), kicker, size=12, color=ORANGE, bold=True)
    txt(slide, Inches(0.55), Inches(0.45), Inches(12), Inches(0.6), title, size=27, color=DARK, bold=True)
    rect(slide, Inches(0.55), Inches(1.12), Inches(12.2), Pt(2), SOFT)


def footer(slide, n):
    txt(slide, Inches(0.55), Inches(7.05), Inches(8), Inches(0.3), "AOIN  ·  Creator Monetization Flow", size=9, color=GREY)
    txt(slide, Inches(11.8), Inches(7.05), Inches(1.0), Inches(0.3), str(n), size=9, color=GREY, align=PP_ALIGN.RIGHT)


# ============================================================ 1. TITLE
s = add_slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(2.55), SW, Inches(0.10), ORANGE)
txt(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.5), "AOIN", size=30, color=ORANGE, bold=True)
txt(s, Inches(1), Inches(2.75), Inches(11.3), Inches(1.4),
    "Creator × Seller\nMonetization Flow", size=46, color=WHITE, bold=True, line_spacing=1.0)
txt(s, Inches(1), Inches(4.9), Inches(11.3), Inches(0.6),
    "An open, self-serve affiliate model — Wishlink-style", size=20, color=ORANGE)
txt(s, Inches(1), Inches(5.6), Inches(11.3), Inches(0.6),
    "Sellers open products for promotion · Creators promote anything · AOIN tracks every sale & pays",
    size=14, color=RGBColor(0xC8, 0xCE, 0xD4))

# ============================================================ 2. THE CORE CONCEPT
s = add_slide(); header(s, "THE BIG IDEA", "The Core Concept — No Matchmaking, Fully Open")
txt(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(0.9),
    "There is no hiring and no selection. Sellers simply make products available for affiliate promotion with a "
    "commission rate. Any creator can promote any available product, get their own link, and earn on the sales they drive.",
    size=16, color=GREY, line_spacing=1.15)
cards = [
    ("Sellers", "Flip products \"on\" for affiliate and set a commission %. That's it — no approving creators.", ORANGE),
    ("Creators", "Browse the open catalog, promote anything via reels or shareable links, earn on every sale.", DARK),
    ("AOIN", "Hosts the catalog, tracks each sale to the right creator, handles payouts, keeps a platform fee.", GREEN),
]
cx = Inches(0.55); cw = Inches(3.95); gap = Inches(0.18)
for i, (t, d, c) in enumerate(cards):
    x = Emu(int(cx) + i * (int(cw) + int(gap)))
    rrect(s, x, Inches(2.55), cw, Inches(2.5), LIGHT, line=SOFT)
    rect(s, x, Inches(2.55), cw, Inches(0.12), c)
    txt(s, Emu(int(x) + Inches(0.3)), Inches(2.85), Emu(int(cw) - Inches(0.6)), Inches(0.5), t, size=20, color=c, bold=True)
    txt(s, Emu(int(x) + Inches(0.3)), Inches(3.45), Emu(int(cw) - Inches(0.6)), Inches(1.5), d, size=14, color=GREY, line_spacing=1.15)
rrect(s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(1.05), SOFT)
txt(s, Inches(0.85), Inches(5.55), Inches(11.6), Inches(0.8),
    "In one line:  Sellers list → Creators self-serve → AOIN attributes every sale to the creator who drove it and pays them.",
    size=16, color=DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 2)

# ============================================================ 3. THE FOUR ROLES
s = add_slide(); header(s, "WHO'S INVOLVED", "The Four Roles")
roles = [
    ("BUYER", "Shops, watches reels, buys. Experience is unchanged — tracking is invisible to them.", DARK),
    ("SELLER (Merchant)", "Lists products; decides which are open for affiliate and at what commission rate.", ORANGE),
    ("CREATOR", "Browses the open catalog, promotes anything via AOIN reels or shareable links, earns commission.", GREEN),
    ("AOIN (Platform)", "Runs the catalog, tracks sales to creators, settles money, takes a platform fee.", GREY),
]
y = Inches(1.5)
for t, d, c in roles:
    rrect(s, Inches(0.55), y, Inches(12.2), Inches(1.2), LIGHT, line=SOFT)
    rect(s, Inches(0.55), y, Inches(0.14), Inches(1.2), c)
    txt(s, Inches(0.95), y, Inches(3.4), Inches(1.2), t, size=18, color=c, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.4), y, Inches(8.1), Inches(1.2), d, size=15, color=GREY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    y = Emu(int(y) + Inches(1.36))
footer(s, 3)

# ============================================================ 4. FLOW 1 — SELLER OPENS PRODUCTS
s = add_slide(); header(s, "FLOW 1 · SELLER", "Seller Makes Products Available for Creators")
bullets(s, Inches(0.55), Inches(1.45), Inches(7.1), Inches(5.2), [
    ("Turn on \"Open for Creators\". ", "On any product, the seller flips a switch to make it promotable."),
    ("Set the commission %. ", "e.g. 10% of each sale — per product, or in bulk across the catalog."),
    ("Optional date window. ", "Run a higher rate during a sale season, then revert."),
    ("Auto-listed instantly. ", "The product appears in the public creator catalog for everyone — no approvals."),
    ("Change or pause anytime. ", "Adjust the rate or remove a product; existing links keep working at the current rate."),
    ("Watch results, not creators. ", "See sales driven, commission owed & paid — without managing individual creators."),
], size=15, gap=12)
rrect(s, Inches(7.95), Inches(1.6), Inches(4.8), Inches(4.7), LIGHT, line=SOFT)
txt(s, Inches(8.2), Inches(1.8), Inches(4.3), Inches(0.5), "The seller's only job", size=15, color=ORANGE, bold=True)
txt(s, Inches(8.2), Inches(2.4), Inches(4.3), Inches(1.2),
    "\"Make it available + set the rate.\"", size=22, color=DARK, bold=True, line_spacing=1.05)
txt(s, Inches(8.2), Inches(3.8), Inches(4.3), Inches(2.3),
    "Everything after that is creator-driven. No inviting, no approving, no negotiating per creator. "
    "Open the product once and it reaches every creator on AOIN.",
    size=14, color=GREY, line_spacing=1.2)
footer(s, 4)

# ============================================================ 5. FLOW 2 — CREATOR ONBOARDING
s = add_slide(); header(s, "FLOW 2 · CREATOR", "Creator Onboarding & Identity")
steps = [
    ("1", "Sign up", "Join as a Creator on the AOIN app (phone OTP)."),
    ("2", "Build identity", "Add Instagram & YouTube handles; upload portfolio videos."),
    ("3", "Pick categories", "Fashion, tech, beauty… so the catalog fits their niche."),
    ("4", "Instant access", "No waiting, no approval — being a creator IS the permission."),
]
x = Inches(0.55); cw = Inches(2.95); gap = Inches(0.17)
for i, (n, t, d) in enumerate(steps):
    xx = Emu(int(x) + i * (int(cw) + int(gap)))
    rrect(s, xx, Inches(2.0), cw, Inches(3.4), LIGHT, line=SOFT)
    circle = s.shapes.add_shape(__import__("pptx").enum.shapes.MSO_SHAPE.OVAL,
                                Emu(int(xx) + Inches(0.3)), Inches(2.3), Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = ORANGE; circle.line.fill.background(); circle.shadow.inherit = False
    tf = circle.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; r = p.add_run(); r.text = n
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
    txt(s, Emu(int(xx) + Inches(0.3)), Inches(3.2), Emu(int(cw) - Inches(0.6)), Inches(0.5), t, size=17, color=DARK, bold=True)
    txt(s, Emu(int(xx) + Inches(0.3)), Inches(3.75), Emu(int(cw) - Inches(0.6)), Inches(1.5), d, size=13.5, color=GREY, line_spacing=1.15)
rrect(s, Inches(0.55), Inches(5.75), Inches(12.2), Inches(0.85), SOFT)
txt(s, Inches(0.85), Inches(5.85), Inches(11.6), Inches(0.65),
    "No application, no review — once set up, the creator instantly has full access to the affiliate catalog.",
    size=15, color=DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)

# ============================================================ 6. FLOW 3 — OPEN CATALOG
s = add_slide(); header(s, "FLOW 3 · CREATOR", "Discovering Products — The Open Catalog")
txt(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(0.7),
    "The heart of the experience. Every product any seller has opened for affiliate is visible to every creator.",
    size=16, color=GREY, line_spacing=1.15)
bullets(s, Inches(0.55), Inches(2.25), Inches(7.0), Inches(4.4), [
    ("One open catalog. ", "All affiliate-enabled products, across all sellers, in one place."),
    ("Each card shows. ", "Product image, price, commission %, category, brand, seller."),
    ("Search & filter. ", "By category, highest commission, brand, price, trending."),
    ("Zero limits. ", "Any product, any number, any seller — creators curate freely."),
], size=15.5, gap=14)
rrect(s, Inches(7.9), Inches(2.25), Inches(4.85), Inches(4.0), WHITE, line=SOFT)
txt(s, Inches(8.15), Inches(2.4), Inches(4.4), Inches(0.4), "Example catalog card", size=12, color=ORANGE, bold=True)
rrect(s, Inches(8.15), Inches(2.85), Inches(4.35), Inches(3.2), LIGHT, line=SOFT)
rect(s, Inches(8.4), Inches(3.1), Inches(3.85), Inches(1.3), SOFT)
txt(s, Inches(8.4), Inches(3.55), Inches(3.85), Inches(0.5), "[ product image ]", size=12, color=GREY, align=PP_ALIGN.CENTER)
txt(s, Inches(8.4), Inches(4.5), Inches(3.85), Inches(0.4), "Blue Running Shoes", size=15, color=DARK, bold=True)
txt(s, Inches(8.4), Inches(4.9), Inches(3.85), Inches(0.4), "₹2,499  ·  SportsBrand", size=13, color=GREY)
pill = rrect(s, Inches(8.4), Inches(5.4), Inches(2.2), Inches(0.5), GREEN)
txt(s, Inches(8.4), Inches(5.4), Inches(2.2), Inches(0.5), "Earn 12% / sale", size=13, color=WHITE, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 6)

# ============================================================ 7. FLOW 4 — TWO WAYS TO PROMOTE
s = add_slide(); header(s, "FLOW 4 · CREATOR", "Promoting a Product — Two Ways")
txt(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(0.6),
    "A creator can use either or both for the same product. No seller-by-seller approval needed.",
    size=15, color=GREY)
# Way A
rrect(s, Inches(0.55), Inches(2.1), Inches(5.95), Inches(2.7), LIGHT, line=SOFT)
rect(s, Inches(0.55), Inches(2.1), Inches(5.95), Inches(0.6), ORANGE)
txt(s, Inches(0.8), Inches(2.15), Inches(5.5), Inches(0.5), "Way A — Make a reel inside AOIN", size=16, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.8), Inches(2.85), Inches(5.5), Inches(1.9), [
    "Tap a product → Create Reel → record/upload",
    "Reel goes live in the AOIN feed, tagged to the creator",
    "Every buyer who watches & buys is credited to them",
], size=14, gap=8)
# Way B
rrect(s, Inches(6.8), Inches(2.1), Inches(5.95), Inches(2.7), LIGHT, line=SOFT)
rect(s, Inches(6.8), Inches(2.1), Inches(5.95), Inches(0.6), GREEN)
txt(s, Inches(7.05), Inches(2.15), Inches(5.5), Inches(0.5), "Way B — Shareable link (\"Wishlink\")", size=16, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(7.05), Inches(2.85), Inches(5.5), Inches(1.9), [
    "Generate a personal link for any product",
    "Share it anywhere — Instagram, YouTube, WhatsApp",
    "Any purchase via the link is credited to them",
], size=14, gap=8)
# Storefront callout
rrect(s, Inches(0.55), Inches(5.05), Inches(12.2), Inches(1.5), SOFT)
txt(s, Inches(0.85), Inches(5.18), Inches(11.6), Inches(0.45), "★  Creator Storefront — the signature feature", size=16, color=ORANGE, bold=True)
txt(s, Inches(0.85), Inches(5.62), Inches(11.6), Inches(0.85),
    "Each creator gets a personal storefront page: one link holding ALL the products they've chosen to promote, "
    "organized by them. They share a single link and their whole audience can browse everything they recommend.",
    size=14, color=DARK, line_spacing=1.15)
footer(s, 7)

# ============================================================ 8. FLOW 5 — BUYER & ATTRIBUTION MOMENT
s = add_slide(); header(s, "FLOW 5 · BUYER", "Buyer Journey & the Moment of Attribution")
chain = [
    "Buyer sees the\nproduct", "Via creator's reel\nOR shared link", "AOIN tags buyer\nto that creator",
    "Buyer browses\n& adds to cart", "Buyer completes\npurchase", "Order stamped with\nthe creator",
]
x = Inches(0.55); bw = Inches(1.85); gap = Inches(0.18); y = Inches(2.4)
for i, c in enumerate(chain):
    xx = Emu(int(x) + i * (int(bw) + int(gap)))
    col = ORANGE if i in (2, 5) else DARK
    rrect(s, xx, y, bw, Inches(1.5), LIGHT if i not in (2,5) else SOFT, line=SOFT)
    txt(s, xx, y, bw, Inches(1.5), c, size=13, color=col, bold=(i in (2,5)),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    if i < len(chain) - 1:
        txt(s, Emu(int(xx) + int(bw) - Inches(0.02)), y, Inches(0.22), Inches(1.5), "›",
            size=22, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rrect(s, Inches(0.55), Inches(4.6), Inches(12.2), Inches(1.6), WHITE, line=SOFT)
txt(s, Inches(0.85), Inches(4.75), Inches(11.6), Inches(0.45), "What the buyer experiences", size=15, color=ORANGE, bold=True)
txt(s, Inches(0.85), Inches(5.2), Inches(11.6), Inches(0.9),
    "A completely normal shopping experience. The \"tag\" connecting the purchase to the creator is invisible to the "
    "buyer — it simply lets AOIN credit the right creator behind the scenes.",
    size=15, color=GREY, line_spacing=1.15)
footer(s, 8)

# ============================================================ 9. FLOW 6 — ATTRIBUTION RULES
s = add_slide(); header(s, "FLOW 6 · RULES", "Attribution — Who Gets Credit")
rules = [
    ("Last-touch wins", "The creator whose reel/link the buyer LAST interacted with before buying gets the credit."),
    ("Attribution window", "The purchase must happen within a set window after the click — e.g. 7 days."),
    ("One sale, one creator", "Each sale is credited to exactly one creator. No splitting, no disputes."),
    ("Returns reverse it", "If the buyer returns/refunds, the credit is removed. Creators earn only on sales that stick."),
]
y = Inches(1.55)
for t, d in rules:
    rrect(s, Inches(0.55), y, Inches(12.2), Inches(1.18), LIGHT, line=SOFT)
    rect(s, Inches(0.55), y, Inches(0.14), Inches(1.18), ORANGE)
    txt(s, Inches(0.95), y, Inches(3.6), Inches(1.18), t, size=17, color=DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.7), y, Inches(7.8), Inches(1.18), d, size=14.5, color=GREY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    y = Emu(int(y) + Inches(1.33))
footer(s, 9)

# ============================================================ 10. FLOW 7 — COMMISSION LIFECYCLE
s = add_slide(); header(s, "FLOW 7 · MONEY", "Commission & Earnings Lifecycle")
stages = [
    ("Order placed", "Pending commission = commission % × sale value", DARK),
    ("AOIN fee deducted", "A small platform fee is taken from the attributed sale", GREY),
    ("PENDING", "Held until the sale is confirmed real", ORANGE),
    ("Delivered + return window passed", "Sale is now genuine and kept", GREEN),
    ("CONFIRMED earnings", "Moves into the creator's wallet", GREEN),
]
y = Inches(1.55);
for i, (t, d, c) in enumerate(stages):
    rrect(s, Inches(0.55), y, Inches(8.6), Inches(0.92), LIGHT, line=SOFT)
    dot = s.shapes.add_shape(__import__("pptx").enum.shapes.MSO_SHAPE.OVAL, Inches(0.8), Emu(int(y)+Inches(0.26)), Inches(0.4), Inches(0.4))
    dot.fill.solid(); dot.fill.fore_color.rgb = c; dot.line.fill.background(); dot.shadow.inherit=False
    tf=dot.text_frame; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=str(i+1)
    r.font.size=Pt(14); r.font.bold=True; r.font.color.rgb=WHITE
    txt(s, Inches(1.45), y, Inches(3.6), Inches(0.92), t, size=15, color=c, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.0), y, Inches(4.0), Inches(0.92), d, size=12.5, color=GREY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    y = Emu(int(y) + Inches(1.02))
# side note
rrect(s, Inches(9.4), Inches(1.55), Inches(3.35), Inches(4.55), SOFT)
txt(s, Inches(9.65), Inches(1.75), Inches(2.9), Inches(0.5), "If cancelled / returned", size=14, color=ORANGE, bold=True)
txt(s, Inches(9.65), Inches(2.35), Inches(2.9), Inches(1.5),
    "before confirmation, the pending commission simply drops off — nothing is paid.",
    size=13.5, color=DARK, line_spacing=1.2)
txt(s, Inches(9.65), Inches(4.0), Inches(2.9), Inches(0.5), "Creator sees, per sale:", size=13, color=ORANGE, bold=True)
txt(s, Inches(9.65), Inches(4.5), Inches(2.9), Inches(1.5),
    "product · sale value · rate · status (pending / confirmed / reversed) · date",
    size=13, color=GREY, line_spacing=1.2)
footer(s, 10)

# ============================================================ 11. FLOW 8 — PAYOUTS
s = add_slide(); header(s, "FLOW 8 · PAYOUTS", "Getting Creators Paid")
steps = [
    ("Earnings accumulate", "Confirmed commissions build up in the creator's balance."),
    ("Minimum threshold", "Withdraw once the balance crosses a minimum (e.g. ₹X)."),
    ("Add payout details", "Creator adds bank / UPI details."),
    ("Regular payout cycle", "AOIN disburses on a schedule (e.g. weekly / monthly batch)."),
    ("Full transparency", "Creator sees payout history; seller sees settlement & their bill."),
]
y = Inches(1.55)
for i,(t,d) in enumerate(steps):
    rrect(s, Inches(0.55), y, Inches(12.2), Inches(0.95), LIGHT, line=SOFT)
    circle = s.shapes.add_shape(__import__("pptx").enum.shapes.MSO_SHAPE.OVAL, Inches(0.8), Emu(int(y)+Inches(0.22)), Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb=ORANGE; circle.line.fill.background(); circle.shadow.inherit=False
    tf=circle.text_frame; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=str(i+1)
    r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=WHITE
    txt(s, Inches(1.6), y, Inches(3.7), Inches(0.95), t, size=16, color=DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.4), y, Inches(7.1), Inches(0.95), d, size=14, color=GREY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    y=Emu(int(y)+Inches(1.07))
footer(s, 11)

# ============================================================ 12. WHAT EACH PARTY SEES
s = add_slide(); header(s, "DASHBOARDS", "What Each Party Sees")
panels = [
    ("SELLER · Web dashboard", ORANGE, [
        "Which products are open + their rates",
        "Sales driven by creators; top creators/reels",
        "Commission owed, paid & upcoming settlement",
        "Toggle products in/out, adjust rates",
    ]),
    ("CREATOR · App + Web", GREEN, [
        "Open catalog (browse / search / filter)",
        "My reels & performance (views, clicks, sales)",
        "My storefront & shareable links",
        "Earnings: pending / confirmed / reversed",
        "Payout balance, withdraw, history",
    ]),
    ("BUYER", DARK, [
        "Nothing new",
        "Normal shopping + reels",
        "Tracking stays invisible",
    ]),
]
x=Inches(0.55); cw=Inches(3.95); gap=Inches(0.18)
for i,(t,c,items) in enumerate(panels):
    xx=Emu(int(x)+i*(int(cw)+int(gap)))
    rrect(s, xx, Inches(1.55), cw, Inches(4.9), LIGHT, line=SOFT)
    rect(s, xx, Inches(1.55), cw, Inches(0.7), c)
    txt(s, Emu(int(xx)+Inches(0.25)), Inches(1.55), Emu(int(cw)-Inches(0.5)), Inches(0.7), t, size=14.5, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, Emu(int(xx)+Inches(0.3)), Inches(2.5), Emu(int(cw)-Inches(0.6)), Inches(3.8), items, size=14, gap=11)
footer(s, 12)

# ============================================================ 13. END-TO-END DIAGRAM
s = add_slide(); header(s, "THE FULL PICTURE", "End-to-End Flow")
flow = [
    ("SELLER opens products for affiliate (sets commission %)", ORANGE),
    ("OPEN CATALOG — every creator sees every available product", DARK),
    ("CREATOR picks any product(s) → makes a reel and/or generates a link / storefront", GREEN),
    ("CREATOR shares — AOIN feed + Instagram / YouTube / anywhere", DARK),
    ("BUYER watches reel or clicks link → arrives tagged to that creator", ORANGE),
    ("BUYER buys within the attribution window → order stamped with creator", DARK),
    ("Delivered + return window passes → commission confirmed (minus AOIN fee)", GREEN),
    ("CREATOR withdraws after minimum · Seller sees settlement · AOIN keeps its fee", ORANGE),
]
y=Inches(1.4); bh=Inches(0.6);
for i,(t,c) in enumerate(flow):
    rrect(s, Inches(1.4), y, Inches(10.5), bh, LIGHT, line=SOFT)
    rect(s, Inches(1.4), y, Inches(0.12), bh, c)
    txt(s, Inches(1.7), y, Inches(10.1), bh, t, size=13.5, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(flow)-1:
        txt(s, Inches(6.5), Emu(int(y)+int(bh)-Inches(0.04)), Inches(0.4), Inches(0.28), "▼", size=12, color=ORANGE, align=PP_ALIGN.CENTER)
    y=Emu(int(y)+int(bh)+Inches(0.11))
footer(s, 13)

# ============================================================ 14. WHY THIS MODEL
s = add_slide(); header(s, "THE PAYOFF", "Why This Model Wins")
wins = [
    ("Open & uncomplicated", "No matchmaking, no approvals, no negotiation. Sellers list, creators self-serve."),
    ("Scales effortlessly", "One opened product reaches every creator; one creator can promote thousands of products."),
    ("Performance-based", "Sellers pay only for sales that actually happen — low risk."),
    ("Creator-empowering", "Creators build their own storefront/brand and earn passively from one shared link."),
    ("Platform flywheel", "More open products → more creators → more reels & links → more buyers → more sellers."),
]
y=Inches(1.55)
for t,d in wins:
    rrect(s, Inches(0.55), y, Inches(12.2), Inches(0.95), LIGHT, line=SOFT)
    rect(s, Inches(0.55), y, Inches(0.14), Inches(0.95), GREEN)
    txt(s, Inches(0.95), y, Inches(3.8), Inches(0.95), t, size=16, color=DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.9), y, Inches(7.6), Inches(0.95), d, size=14, color=GREY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    y=Emu(int(y)+Inches(1.06))
footer(s, 14)

# ============================================================ 15. CLOSING
s = add_slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(3.1), SW, Inches(0.08), ORANGE)
txt(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.0), "In one line", size=18, color=ORANGE, bold=True)
txt(s, Inches(1), Inches(3.3), Inches(11.3), Inches(2.6),
    "Sellers flip products \"on\" for affiliate with a commission rate; every creator can freely promote any of them "
    "via AOIN reels or shareable links; AOIN tracks each sale to the creator who drove it and pays them — keeping a "
    "platform fee. Fully open, self-serve, Wishlink-style.",
    size=24, color=WHITE, bold=True, line_spacing=1.15)

prs.save("/Users/nikhilpatel/Projects/aoin/Ecommerce_Backend/docs/AOIN_Creator_Seller_Flow.pptx")
print("SAVED")
